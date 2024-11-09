
from pptx import Presentation
from pptx.util import Inches
import markdown
from bs4 import BeautifulSoup

def markdown_to_html(markdown_text):
    return markdown.markdown(markdown_text)

def html_to_text(html):
    soup = BeautifulSoup(html, 'html.parser')
    return soup.get_text()

def add_slide(prs, title, content):
    slide_layout = prs.slide_layouts[1]  # Use the title and content layout
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]

    title_placeholder.text = title
    content_placeholder.text = content

def main():
    prs = Presentation()
    with open('web3-blockchain.md', 'r', encoding='utf-8') as file:
        markdown_text = file.read()

    html = markdown_to_html(markdown_text)
    soup = BeautifulSoup(html, 'html.parser')

    current_title = None
    current_content = []

    for element in soup:
        if element.name == 'h1':
            if current_title:
                add_slide(prs, current_title, '\n'.join(current_content))
            current_title = element.get_text()
            current_content = []
        elif element.name in ['h2', 'h3', 'h4', 'h5', 'h6']:
            if current_title:
                add_slide(prs, current_title, '\n'.join(current_content))
            current_title = element.get_text()
            current_content = []
        else:
            current_content.append(element.get_text())

    if current_title:
        add_slide(prs, current_title, '\n'.join(current_content))

    prs.save('web3-blockchain.pptx')

if __name__ == "__main__":
    main()